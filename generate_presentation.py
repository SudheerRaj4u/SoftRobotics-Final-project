from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Theme Configuration (Tech/Dark Mode) ---
BG_COLOR = RGBColor(15, 23, 42)      # Deep Navy/Black
ACCENT_1 = RGBColor(56, 189, 248)    # Cyan/Sky Blue
ACCENT_2 = RGBColor(168, 85, 247)    # Purple/Neon
TEXT_MAIN = RGBColor(241, 245, 249)  # White/Off-White
TEXT_SUB = RGBColor(148, 163, 184)   # Greyish

# --- Image Paths ---
# IMPORTANT: Adjust these filenames to match the actual generated files
# Since the system generates random timestamps, I will look for the most recent matching files in the directory
# For robustness, we will try to find them dynamically or assume they are renamed.
# Since I cannot easily rename in this script without knowing the exact timestamp,
# I will assume the files are in the artifacts directory and search for them.
ARTIFACT_DIR = r"C:\Users\sudhe\.gemini\antigravity\brain\9b83bb05-30b8-472b-9a11-87383433278a"

def get_image_path(basename_prefix):
    # Find file starting with basename_prefix in ARTIFACT_DIR
    try:
        files = os.listdir(ARTIFACT_DIR)
        for f in files:
            if f.startswith(basename_prefix) and f.endswith(".png"):
                return os.path.join(ARTIFACT_DIR, f)
    except Exception as e:
        print(f"Error finding image {basename_prefix}: {e}")
    return None

def create_presentation():
    prs = Presentation()
    
    # --- Helper Functions ---
    def create_slide(layout_idx=6): 
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return slide

    def add_bullets(slide, items, left=Inches(0.5), top=Inches(1.5), width=Inches(9)):
        txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = "Arial"
            p.space_after = Pt(14)

    def add_image(slide, img_path, left, top, height=None, width=None):
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height, width=width)

    # --- Slide 1: Title ---
    slide = create_slide()
    
    # Image Background (Covering Right Half or faint overlay)
    img_path = get_image_path("soft_robot_futuristic")
    if img_path:
        # Add stylistic full height image on right
        add_image(slide, img_path, Inches(5), 0, height=prs.slide_height)
    
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(3))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Development of a\nComprehensive Simulation Framework"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_1
    
    p2 = tf.add_paragraph()
    p2.text = "Multi-Channel Soft Continuum Robots"
    p2.font.size = Pt(28)
    p2.font.color.rgb = ACCENT_2
    p2.space_before = Pt(20)

    # --- Slide 2: Physics Engine (Novelty 1) ---
    slide = create_slide()
    # Header
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tbox.text_frame.text = "NOVELTY I: PHYSICS ENGINE"
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_2
    tbox.text_frame.paragraphs[0].font.active_color = ACCENT_2
    tbox.text_frame.paragraphs[0].font.size = Pt(32)
    tbox.text_frame.paragraphs[0].font.bold = True

    add_bullets(slide, [
        "Continuum Mechanics Chain (5-Step):",
        "1. Pressure -> Axial Force",
        "2. Asymmetry -> Moment Calculation",
        "3. Constitutive Law (Neo-Hookean)",
        "4. Arc Geometry Derivation",
        "5. Kinematic Transformation"
    ], width=Inches(5))

    img_path = get_image_path("physics_wireframe_specific")
    if img_path:
        add_image(slide, img_path, Inches(5.5), Inches(1.5), height=Inches(4.5))

    # --- Slide 3: 4-Channel Actuation (Novelty 2) ---
    slide = create_slide()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tbox.text_frame.text = "NOVELTY II: 4-CHANNEL ACTUATION"
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_2
    tbox.text_frame.paragraphs[0].font.size = Pt(32)
    tbox.text_frame.paragraphs[0].font.bold = True

    add_bullets(slide, [
        "Beyond Planar Bending:",
        "- Independent control of 4 chambers.",
        "- Enables arbitrary Bending Plane (phi).",
        "- Allows Stiffness Control via co-contraction."
    ], width=Inches(5))

    img_path = get_image_path("cross_section_4ch")
    if img_path:
        add_image(slide, img_path, Inches(5.5), Inches(1.5), height=Inches(4.5))

    # --- Slide 4: Advanced Analysis (Novelty 3) ---
    slide = create_slide()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tbox.text_frame.text = "NOVELTY III: ADVANCED ANALYSIS"
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_2
    tbox.text_frame.paragraphs[0].font.size = Pt(32)
    tbox.text_frame.paragraphs[0].font.bold = True

    add_bullets(slide, [
        "Monte Carlo Workspace Visualization:",
        "- Sampling 5000+ random pressure vectors.",
        "- Generating 3D point cloud of reachable space.",
        "Parametric Sweep Analysis:",
        "- Sweeping L, tw, P to find optimal designs."
    ], width=Inches(5))

    img_path = get_image_path("workspace_viz")
    if img_path:
        add_image(slide, img_path, Inches(5.5), Inches(1.5), height=Inches(4.5))

    # --- Slide 5: Technical - Material Models ---
    slide = create_slide()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tbox.text_frame.text = "TECHNICAL DEPTH: MATERIAL MODELS"
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
    tbox.text_frame.paragraphs[0].font.size = Pt(32)
    tbox.text_frame.paragraphs[0].font.bold = True

    add_bullets(slide, [
        "Hyperelastic Constitutive Laws:",
        "- Neo-Hookean: Strain Energy based on 1st Invariant (I1).",
        "- Mooney-Rivlin: Captures higher-order non-linearities.",
        "- Ogden: Phenomenological model for large strains.",
        "Significance: Matches real silicone (EcoFlex 00-30) behavior."
    ])

    # --- Slide 6: Technical - Inverse Kinematics ---
    slide = create_slide()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tbox.text_frame.text = "TECHNICAL DEPTH: INVERSE KINEMATICS"
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
    tbox.text_frame.paragraphs[0].font.size = Pt(32)
    tbox.text_frame.paragraphs[0].font.bold = True

    add_bullets(slide, [
        "Hybrid Solver Architecture:",
        "- Global Search: Differential Evolution (Avoids Local Minima).",
        "- Local Refinement: L-BFGS-B (Precision < 10^-5m).",
        "Robustness:",
        "- Validated across 2 to 100 segment scales.",
        "- Solves for Pressure Vector P given Target T."
    ])

    # --- Slide 5: Conclusion ---
    slide = create_slide()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tbox.text_frame.text = "CONCLUSION"
    tbox.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
    tbox.text_frame.paragraphs[0].font.size = Pt(36)
    
    add_bullets(slide, [
        "Summary of Contributions:",
        "- [x] Novel 4-Channel Physics Engine.",
        "- [x] Rigorous Validation Chain.",
        "- [x] Advanced Design Tools (Sweeps/Space).",
        "Significance: A robust framework for soft robot design."
    ])

    # Save
    filename = "Soft_Robotics_Presentation_Keynote.pptx"
    prs.save(filename)
    print(f"Presentation saved to {filename}")

if __name__ == "__main__":
    create_presentation()
